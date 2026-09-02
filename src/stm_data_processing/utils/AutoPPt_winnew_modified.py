#!/Users/hunfen/Documents/GitHub/STM_DataProcessing/.venv/bin/python

# fmt: off
import shutil
import sys
from datetime import datetime
from pathlib import Path

import cv2
import matplotlib
import matplotlib.pyplot as plt
import nanonispy as nap
import numpy as np
from matplotlib.animation import FuncAnimation
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches, Pt

np.float = float
np.int = int
# fmt: on

# isort: off
try:
    from .plot_funcs import (
        build_bias_labels,
        finite_range,
        subtractMeanPlane,
        topo_colormap,
    )
except ImportError:  # allow running as a standalone script
    from plot_funcs import (
        build_bias_labels,
        finite_range,
        subtractMeanPlane,
        topo_colormap,
    )
# isort: on

# ================= 交互式输入区 =================
if __name__ == "__main__":
    # 1. 数据文件夹路径
    print("请输入数据处理文件夹路径(多个路径用英文逗号分隔):")
    input_paths = input("> ").strip()
    if input_paths:
        DataFolderpath = [Path(p.strip()) for p in input_paths.split(",")]
    else:
        sys.exit("错误:未输入任何数据文件夹路径,脚本终止。")

    # 2. 输出存储路径
    print("\n请输入PPT存放目录路径(直接回车则保存在数据文件夹路径下):")
    input_storage = input("> ").strip()
    if input_storage:
        Storagepath = Path(input_storage)
    else:
        # 默认保存到第一个数据文件夹中
        Storagepath = Path(DataFolderpath[0])
        print(f"未指定输出路径,将PPT保存在数据文件夹:{Storagepath}")

    # 3. PPT文件名
    print("\n请输入PPT文件名(不含路径,直接回车自动生成):")
    input_pptname = input("> ").strip()
    if input_pptname:
        if not input_pptname.endswith(".pptx"):
            input_pptname += ".pptx"
        PPTname = input_pptname
    else:
        # 自动以日期命名
        from datetime import datetime

        PPTname = f"auto_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        print(f"自动生成文件名: {PPTname}")

    # 4. QPI开关(on/off)
    print("\n是否输出QPI图?输入 on 或 off(直接回车默认 on):")
    input_qpi = input("> ").strip().lower()
    QPI_switch = "on" if input_qpi in ("", "on") else "off"
    print(f"QPI输出: {QPI_switch}")

    # 5. map电流开关(on/off)
    print("\n是否输出Map中的电流图?输入 on 或 off(直接回车默认 off):")
    input_mapi = input("> ").strip().lower()
    mapI_switch = "on" if input_mapi == "on" else "off"
    print(f"Map电流输出: {mapI_switch}")


prs = Presentation()
matplotlib.use("Agg")
plt.rcParams["font.sans-serif"] = ["Arial"]
plt.rcParams["axes.unicode_minus"] = True
plt.rcParams.update({"font.size": 22})

DatFiles = []
SxmFiles = []
LinecutFiles = []
MapFiles = []
Files = []
global slide


# 读取目标文件夹中的所有文件的名字并储存在Files中
def ReadAllFiles(folder_path: Path):
    for file in folder_path.rglob("*"):
        if file.is_file():
            Files.append(file)


# 按照不同的数据格式给文件分类
def SortFiles(files):
    for each in files:
        if each.suffix == ".dat":
            # DatFiles.append(each)
            try:
                nap.read.Spec(each).signals["Bias calc (V)"]
                DatFiles.append(each)
            except Exception:
                pass
        if each.suffix == ".sxm":
            SxmFiles.append(each)
        if each.suffix == ".3ds":
            raw_data_1 = nap.read.Grid(each)
            try:
                data_1 = raw_data_1.signals["LI Demod 1 X (A)"]
            except Exception:
                data_1 = raw_data_1.signals["LI Demod 1 X [AVG] (A)"]
            if np.shape(data_1)[0] == 1:
                LinecutFiles.append(each)
            if np.shape(data_1)[0] != 1:
                MapFiles.append(each)


def get_creation_time(file_path):
    stat = Path(file_path).stat()
    # Prefer the true creation/birth time where the platform exposes it
    # (macOS, FreeBSD, Windows). Linux has no portable birth time via
    # os.stat: st_ctime is the inode-change time (not the creation time),
    # so the last-modified time is used as the documented fallback proxy.
    if hasattr(stat, "st_birthtime") and stat.st_birthtime > 0:
        return stat.st_birthtime
    return stat.st_mtime


def sort_files_by_creation_time(files):
    # 获取每个文件的创建时间
    file_creation_times = [(file, get_creation_time(file)) for file in files]
    # 按照创建时间排序
    sorted_files = sorted(file_creation_times, key=lambda x: x[1])
    # 提取排序后的文件路径
    sorted_file_paths = [file[0] for file in sorted_files]
    return sorted_file_paths


# 旋转矩阵的函数,其中要把header文件中的scan_angle传递进去
def img_rotate(data, angle, range_x, range_y):
    # 获取图像中心点坐标以及行列数
    image = data
    image = image.astype(np.float32)
    rows, cols = image.shape
    center = ((cols - 1) / 2.0, (rows - 1) / 2.0)

    # 定义旋转矩阵
    rotation_matrix = cv2.getRotationMatrix2D(center, angle, 1)

    # 计算旋转后的图像大小
    cos_theta = np.abs(rotation_matrix[0, 0])
    sin_theta = np.abs(rotation_matrix[0, 1])
    cols_new = int(rows * sin_theta + cols * cos_theta)
    rows_new = int(rows * cos_theta + cols * sin_theta)

    # 计算平移矩阵,确保旋转后图像完整显示在新的图像中
    rotation_matrix[0, 2] += (cols_new - cols) / 2
    rotation_matrix[1, 2] += (rows_new - rows) / 2

    # 进行仿射变换
    rotated_image = cv2.warpAffine(
        image, rotation_matrix, (cols_new, rows_new), flags=cv2.INTER_NEAREST
    )  # borderValue=(255,255,255
    # 计算新的图的坐标尺寸
    range_x_new = range_y * sin_theta + range_x * cos_theta
    range_y_new = range_y * cos_theta + range_x * sin_theta
    return rotated_image, range_x_new, range_y_new


def img_rotate_for_box(data, degree=90, zoom_pan=1):
    img = data
    rows, cols = img.shape
    M = cv2.getRotationMatrix2D(((cols - 1) / 2.0, (rows - 1) / 2.0), degree, zoom_pan)
    dst = cv2.warpAffine(img, M, (cols, rows))
    return dst


# 绘制Map的函数
def ShowMap(mappath, n):
    divider = 1
    if "d1" in str(mappath):
        divider = 1
    if "d10" in str(mappath):
        divider = 10
    if "d100" in str(mappath):
        divider = 100
    raw_data = nap.read.Grid(mappath)
    bias = build_bias_labels(raw_data, divider)
    try:
        data = raw_data.signals["LI Demod 1 X (A)"][:, :, n]
    except Exception:
        data = raw_data.signals["LI Demod 1 X [AVG] (A)"][:, :, n]
    # 下面这行 angle 在函数内未使用,可删除
    # angle = float(raw_data.header["angle"])
    scan_range = raw_data.header["size_xy"]

    # 绘制map
    fig = plt.figure(figsize=(2.55, 2.55))
    ax0 = fig.add_subplot(111)
    ax0.imshow(
        data,
        origin="lower",
        extent=(0, scan_range[0] * 1e9, 0, scan_range[1] * 1e9),
        cmap=topo_colormap("rainbow"),
    )
    plt.xticks([])
    plt.yticks([])
    ax0.text(
        0.04,
        0.90,
        f"{bias[n]:.2f} mV",
        transform=ax0.transAxes,
        fontdict={"family": "Arial", "size": "13", "color": "black", "weight": "bold"},
    )
    plt.axis("off")
    # 如果 Storagepath 已转为 Path 对象,请同步修改保存路径为:
    # plt.savefig(str(Storagepath / "folder_map") / f"temp_map_{n}.tif", ...)
    plt.savefig(
        str(Storagepath / "folder_map" / f"temp_map_{n}.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()


def QPI(mappath, n):
    divider = 1
    if "d1" in str(mappath):
        divider = 1
    if "d10" in str(mappath):
        divider = 10
    if "d100" in str(mappath):
        divider = 100
    raw_data = nap.read.Grid(mappath)
    bias = build_bias_labels(raw_data, divider)
    try:
        data = raw_data.signals["LI Demod 1 X (A)"][:, :, n]
    except Exception:
        data = raw_data.signals["LI Demod 1 X [AVG] (A)"][:, :, n]
    scan_range = raw_data.header["size_xy"]
    # 绘制QPI,先做完fft,再转动
    fft2 = np.fft.fft2(data)
    shift2center = np.fft.fftshift(fft2)
    QPI2 = np.log(1 + abs(shift2center))
    # 由于colorbar范围问题,需要求平均值和标准差以将colorbar颜色卡在某个范围
    if np.any(np.isfinite(QPI2)):
        mean_value = float(np.nanmean(QPI2))
        std_deviation = float(np.nanstd(QPI2))
        vqmin = float(np.nanmin(QPI2))
        vqmax = mean_value + 1.5 * std_deviation
    else:
        vqmin, vqmax = None, None
    # 计算Q空间的范围
    range_qx = 2 * np.pi / (scan_range[0] * 1e9 / len(data[0]))
    range_qy = 2 * np.pi / (scan_range[1] * 1e9 / len(data[:, 0]))
    fig = plt.figure(figsize=(2.55, 2.55))
    ax0 = fig.add_subplot(111)
    ax0.imshow(
        QPI2,
        origin="lower",
        cmap=topo_colormap("gray_r"),
        vmin=vqmin,
        vmax=vqmax,
        extent=(-range_qx / 2, range_qx / 2, -range_qy / 2, range_qy / 2),
    )
    ax0.text(
        0.04,
        0.90,
        f"{bias[n]:.2f} mV",
        transform=ax0.transAxes,
        fontdict={"family": "Arial", "size": "13", "color": "red", "weight": "bold"},
    )
    plt.xticks([])
    plt.yticks([])
    plt.axis("off")
    plt.savefig(
        str(Storagepath / "folder_QPI" / f"temp_QPI_{n}.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()


def ShowMapI(mappath, n):
    divider = 1
    if "d1" in str(mappath):
        divider = 1
    if "d10" in str(mappath):
        divider = 10
    if "d100" in str(mappath):
        divider = 100
    raw_data = nap.read.Grid(mappath)
    # 获取mapping对应的偏压列表
    bias = build_bias_labels(raw_data, divider)
    data = raw_data.signals["Current (A)"][:, :, n]
    scan_range = raw_data.header["size_xy"]

    # --------------------------------------------------------------------------------
    # 绘制电流map
    fig = plt.figure(figsize=(2.55, 2.55))
    ax0 = fig.add_subplot(111)
    ax0.imshow(
        data,
        origin="lower",
        extent=(0, scan_range[0] * 1e9, 0, scan_range[1] * 1e9),
        cmap=topo_colormap("rainbow"),
    )
    plt.xticks([])
    # plt.tick_params(axis='x', which='both', pad=0)
    plt.yticks([])
    # plt.tick_params(axis='y', which='both', pad=0)
    ax0.text(
        0.04,
        0.90,
        f"{bias[n]:.2f} mV",
        transform=ax0.transAxes,
        fontdict={"family": "Arial", "size": "13", "color": "black", "weight": "bold"},
    )
    # ax0.set_xlabel('X (nm)',fontdict={'family':'Arial','size':'10'})
    # ax0.set_ylabel('Y (nm)',fontdict={'family':'Arial','size':'10'})
    # cax = fig.add_axes([ax0.get_position().x1+0.01,ax0.get_position().y0,0.02,ax0.get_position().height])  #添加一个colorbar的位置
    # plt.colorbar(map1,cax=cax)
    plt.axis("off")
    plt.savefig(
        str(Storagepath / "folder_mapI" / f"temp_mapI_{n}.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()


# 绘制形貌图的函数
def SXM(topopath):
    try:
        raw_data = nap.read.Scan(topopath)
        topo = raw_data.signals["Z"]["forward"]
        size = raw_data.header["scan_range"]
        direction = raw_data.header["scan_dir"]
        fig = plt.figure(figsize=(2.55, 2.55))
        fig.add_subplot(111)
        # 这里扫图方向不仅影响着储存在矩阵中的行的正反,还关系着转动矩阵转角符号问题
        if direction == "up":
            plt.imshow(
                topo,
                origin="lower",
                cmap=topo_colormap("Blues_r"),
                extent=(0, size[0] * 1e9, 0, size[1] * 1e9),
            )
        else:
            plt.imshow(
                topo,
                cmap=topo_colormap("Blues_r"),
                extent=(0, size[0] * 1e9, 0, size[1] * 1e9),
            )
        plt.xticks([])
        plt.yticks([])
        plt.axis("off")
        plt.savefig(
            Storagepath / "temp.tif",
            format="tif",
            bbox_inches="tight",
            transparent=True,
            pad_inches=0,
        )
        plt.close(fig)
    except Exception as e:
        raise ValueError(f"Failed to read SXM file {topopath}: {e}") from e


# 寻找单谱或linecut数据对应的topography文件
def find_nearest_file(target_file_path, file_paths):
    target_creation_time = get_creation_time(target_file_path)
    nearest_file_path = None
    min_time_difference = float("inf")

    for file_path in file_paths:
        try:
            current_creation_time = get_creation_time(file_path)
            time_difference = target_creation_time - current_creation_time

            if time_difference < min_time_difference and time_difference > 0:
                min_time_difference = time_difference
                nearest_file_path = file_path
        except Exception:
            # Handle the case where the file does not exist
            pass
    if nearest_file_path is None:
        nearest_file_path = "Topography Not Found"
    return nearest_file_path


# 绘制单谱的函数,并且搜索到最近的形貌图在形貌图上标注位置,保存单谱图和对应的形貌图
def STS(stspath):
    topopath = find_nearest_file(stspath, SxmFiles)
    raw_data = nap.read.Spec(stspath)
    divider = 1
    if "d1" in str(stspath):
        divider = 1
    if "d10" in str(stspath):
        divider = 10
    if "d100" in str(stspath):
        divider = 100
    bias = raw_data.signals["Bias calc (V)"] * 1000 / divider
    try:
        didv = raw_data.signals["LI Demod 1 X [AVG] (A)"]
        # current = raw_data.signals['Current [AVG] (A)']
    except Exception:
        didv = raw_data.signals["LI Demod 1 X (A)"]
        # current = raw_data.signals['Current (A)']
    # 搜索作单谱位置的坐标
    X, Y = eval(raw_data.header["X (m)"]) * 1e9, eval(raw_data.header["Y (m)"]) * 1e9
    # 生成单谱图并保存

    fig = plt.figure(figsize=(2.15, 2.15))
    ax1 = fig.add_subplot(111)
    ax1.tick_params(axis="both", which="major", pad=1)
    matplotlib.rcParams["font.size"] = 6
    ax1.plot(bias, didv, "r-")
    font = {"family": "Arial", "size": 6}
    ax1.set_xlabel("Bias (mV)", fontdict=font, labelpad=0.1)
    ax1.set_ylabel(r"d$\it{I}$/d$\it{V}$ (a.u.)", fontdict=font, labelpad=0.1)
    plt.savefig(
        str(Storagepath / "temp_sts.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()

    if topopath != "Topography Not Found":
        # 生成形貌并画出单谱位置,生成图片
        raw_data = nap.read.Scan(topopath)
        topo = raw_data.signals["Z"]["forward"]
        scan_range = raw_data.header["scan_range"]
        scan_offset = raw_data.header["scan_offset"]
        angle = float(raw_data.header["scan_angle"])
        direction = raw_data.header["scan_dir"]
        fig = plt.figure(figsize=(2.55, 2.55))
        ax0 = fig.add_subplot(111)
        if direction == "up":
            angle = angle
            topo2, range_x_new, range_y_new = img_rotate(
                topo, angle, scan_range[0], scan_range[1]
            )
            ax0.imshow(
                topo2,
                origin="lower",
                cmap=topo_colormap("Blues_r"),
                vmin=finite_range(topo)[0],
                vmax=finite_range(topo)[1],
                extent=(
                    (scan_offset[0] - range_x_new / 2) * 1e9,
                    (scan_offset[0] + range_x_new / 2) * 1e9,
                    (scan_offset[1] - range_y_new / 2) * 1e9,
                    (scan_offset[1] + range_y_new / 2) * 1e9,
                ),
            )
        else:
            angle = -angle
            topo2, range_x_new, range_y_new = img_rotate(
                topo, angle, scan_range[0], scan_range[1]
            )
            ax0.imshow(
                topo2,
                cmap=topo_colormap("Blues_r"),
                vmin=finite_range(topo)[0],
                vmax=finite_range(topo)[1],
                extent=(
                    (scan_offset[0] - range_x_new / 2) * 1e9,
                    (scan_offset[0] + range_x_new / 2) * 1e9,
                    (scan_offset[1] - range_y_new / 2) * 1e9,
                    (scan_offset[1] + range_y_new / 2) * 1e9,
                ),
            )
        ax0.plot(X, Y, "ro", markersize=3)
        # if (extentx[0])<=X<=extentx[1] and extenty[0]<=Y<=extenty[1]:
        #     ax0.plot(X,Y,'ro',markersize=3)
        # else:
        #     ax0.text(0.2,0.5,'Location Not Found',transform=ax0.transAxes,fontdict={'family':'Arial','size':'10','color':'red','weight': 'bold'})

        plt.xticks([])
        plt.yticks([])
        plt.axis("off")
        plt.savefig(
            str(Storagepath / "temp_ststopo.tif"),
            format="tif",
            bbox_inches="tight",
            transparent=True,
            pad_inches=0,
        )
        plt.close()


# 绘制linecut的函数,并且搜索到最近的形貌图在形貌图上标注位置,保存linecut图和对应的形貌图
def Linecut(lcpath):
    topopath = find_nearest_file(lcpath, SxmFiles)
    raw_data = nap.read.Grid(lcpath)
    L = raw_data.header["size_xy"][0] * 1e9  # 获取linecut长度,单位为nm
    divider = 1
    if "d1" in str(lcpath):
        divider = 1
    if "d10" in str(lcpath):
        divider = 10
    if "d100" in str(lcpath):
        divider = 100
    bias = raw_data.signals["sweep_signal"][:] * 1000 / divider
    try:
        lcdata = raw_data.signals["LI Demod 1 X [AVG] (A)"][0, :, :]
    except Exception:
        lcdata = raw_data.signals["LI Demod 1 X (A)"][0, :, :]

    # 画出linecut的瀑布图
    fig, (ax0, ax1, ax2) = plt.subplots(
        1, 3, figsize=(5.8, 3), gridspec_kw={"width_ratios": [2, 2, 1]}
    )
    ax1.imshow(
        lcdata,
        extent=(bias[0], bias[-1], 0, L),
        aspect="auto",
        origin="lower",
        cmap=topo_colormap("rainbow"),
        interpolation="none",
    )
    ax1.tick_params(axis="both", which="major", pad=1)
    matplotlib.rcParams["font.size"] = 6
    font = {"family": "Arial", "size": 8}
    ax1.set_xlabel("Bias (mV)", fontdict=font, labelpad=0.3)
    ax1.set_ylabel("Length(nm)", fontdict=font, labelpad=0.3)

    # 画出linecut的堆叠偏移单谱
    offset = 0.25 * np.nanmean(lcdata[0])
    cmap = plt.get_cmap("brg")
    colors = [
        cmap(i) for i in np.linspace(0, 1, len(lcdata))
    ]  # 获得渐变色以绘制linecut
    for i in range(len(lcdata)):
        spec = lcdata[i] + 1 * i * offset
        ax0.plot(bias, spec, linewidth=0.2, color=colors[i])
    ax0.axhline(y=0, xmin=0.3, xmax=0.7, c="k", lw=0.1, ls="-")

    ax0.text(bias[-1], spec[-1], len(lcdata))
    ax0.text(bias[-1], lcdata[0][-1], 0)

    ax0.tick_params(axis="both", which="major", pad=1)
    matplotlib.rcParams["font.size"] = 6
    plt.yticks([])
    font = {"family": "Arial", "size": 8}
    ax0.set_xlabel("Bias (mV)", fontdict=font, labelpad=0.3)
    ax0.set_ylabel(r"d$\it{I}$/d$\it{V}$ (a.u.)", fontdict=font, labelpad=0.1)

    # 画出linecut的形貌线
    ax2.tick_params(axis="both", which="major", pad=1)
    matplotlib.rcParams["font.size"] = 10
    font = {"family": "Arial", "size": 8}
    ax2.set_xlabel("Height (pm)", fontdict=font, labelpad=0.3)
    height = raw_data.signals["topo"][0][:] * 1e12
    ax2.plot(height - min(height), np.linspace(0, L, len(height)))
    ax2.set_ylim(0, L)

    plt.subplots_adjust(wspace=0.25)  # wspace 控制水平间距
    plt.savefig(
        str(Storagepath / "temp_lc.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()

    # 画出overlap图
    cmap = plt.get_cmap("brg")
    colors = [
        cmap(i) for i in np.linspace(0, 1, len(lcdata))
    ]  # 获得渐变色以绘制linecut
    fig = plt.figure(figsize=(2.8, 2.8))
    ax0 = fig.add_subplot(111)
    for i in range(len(lcdata)):
        ax0.plot(bias, lcdata[i], linewidth=0.5, color=colors[i])
    ax0.axhline(y=0, xmin=0.3, xmax=0.7, c="k", lw=0.1, ls="-")
    ax0.tick_params(axis="both", which="major", pad=1)
    matplotlib.rcParams["font.size"] = 6
    font = {"family": "Arial", "size": 8}
    ax0.set_xlabel("Bias (mV)", fontdict=font, labelpad=0.3)
    ax0.set_ylabel(r"d$\it{I}$/d$\it{V}$ (a.u.)", fontdict=font, labelpad=0.1)
    plt.savefig(
        str(Storagepath / "temp_ol.tif"),
        format="tif",
        bbox_inches="tight",
        transparent=True,
        pad_inches=0,
    )
    plt.close()
    # 画出对应的形貌图
    if topopath != "Topography Not Found":
        raw_data = nap.read.Scan(topopath)
        topo = raw_data.signals["Z"]["forward"]
        scan_range = raw_data.header["scan_range"]
        scan_offset = raw_data.header["scan_offset"]
        angle = float(raw_data.header["scan_angle"])
        direction = raw_data.header["scan_dir"]

        linecut_raw_data = nap.read.Grid(lcpath)
        linecut_size = linecut_raw_data.header["size_xy"][0]
        linecut_center = linecut_raw_data.header["pos_xy"]
        lcangle = linecut_raw_data.header["angle"] * np.pi / 180
        X = [
            (linecut_center[0] - linecut_size / 2 * np.cos(lcangle)) * 1e9,
            (linecut_center[1] + linecut_size / 2 * np.sin(lcangle)) * 1e9,
        ]
        Y = [2 * linecut_center[0] * 1e9 - X[0], 2 * linecut_center[1] * 1e9 - X[1]]
        fig = plt.figure(figsize=(2.8, 2.8))
        ax0 = fig.add_subplot(111)
        if direction == "up":
            angle = angle
            topo2, range_x_new, range_y_new = img_rotate(
                topo, angle, scan_range[0], scan_range[1]
            )
            ax0.imshow(
                topo2,
                origin="lower",
                cmap=topo_colormap("Blues_r"),
                vmin=finite_range(topo)[0],
                vmax=finite_range(topo)[1],
                extent=(
                    (scan_offset[0] - range_x_new / 2) * 1e9,
                    (scan_offset[0] + range_x_new / 2) * 1e9,
                    (scan_offset[1] - range_y_new / 2) * 1e9,
                    (scan_offset[1] + range_y_new / 2) * 1e9,
                ),
            )

        else:
            angle = -angle
            topo2, range_x_new, range_y_new = img_rotate(
                topo, angle, scan_range[0], scan_range[1]
            )
            ax0.imshow(
                topo2,
                cmap=topo_colormap("Blues_r"),
                vmin=finite_range(topo)[0],
                vmax=finite_range(topo)[1],
                extent=(
                    (scan_offset[0] - range_x_new / 2) * 1e9,
                    (scan_offset[0] + range_x_new / 2) * 1e9,
                    (scan_offset[1] - range_y_new / 2) * 1e9,
                    (scan_offset[1] + range_y_new / 2) * 1e9,
                ),
            )
        # if (extentx[0]<=X[0]<=extentx[1] and extenty[0]<=X[1]<=extenty[1]) or (extentx[0]<=Y[0]<=extentx[1] and extenty[0]<=Y[1]<=extenty[1]):
        ax0.plot(X[0], X[1], color="k", fillstyle="none", marker="o", markersize=5)
        ax0.plot([X[0], Y[0]], [X[1], Y[1]], "k--")
        # else:
        #     ax0.text(0.2,0.5,'Location Not Found',transform=ax0.transAxes,fontdict={'family':'Arial','size':'10','color':'red','weight': 'bold'})
        plt.xticks([])
        plt.yticks([])
        plt.axis("off")
        plt.savefig(
            str(Storagepath / "temp_lctopo.tif"),
            format="tif",
            bbox_inches="tight",
            transparent=True,
            pad_inches=0,
        )
        plt.close()


# 创建不同偏压map动图的函数
def update_map(n):
    ax.cla()  # 清除当前轴上的内容
    ax.set_axis_off()
    img_path = str(Storagepath / "folder_map") / f"temp_map_{n}.tif"
    img = plt.imread(img_path)
    # 获取图像的实际范围
    range_x, range_y = img.shape[1], img.shape[0]
    # 使用 extent 参数设置图像的实际范围
    ax.imshow(img, extent=[0, range_x, 0, range_y])


# 创建不同偏压QPI动图的函数
def update_QPI(n):
    ax.cla()  # 清除当前轴上的内容
    ax.set_axis_off()
    img_path = str(Storagepath / "folder_QPI") / f"temp_QPI_{n}.tif"
    img = plt.imread(img_path)
    # 获取图像的实际范围
    range_x, range_y = img.shape[1], img.shape[0]
    # 使用 extent 参数设置图像的实际范围
    ax.imshow(img, extent=[0, range_x, 0, range_y])


# 创建不同偏压电流动图函数
def update_mapI(n):
    ax.cla()  # 清除当前轴上的内容
    ax.set_axis_off()
    img_path = str(Storagepath / "folder_mapI") / f"temp_mapI_{n}.tif"
    img = plt.imread(img_path)
    # 获取图像的实际范围
    range_x, range_y = img.shape[1], img.shape[0]
    # 使用 extent 参数设置图像的实际范围
    ax.imshow(img, extent=[0, range_x, 0, range_y])


# 创建角度识别函数,要求角度是恒正的
def angle_def(angle):
    if angle >= 0:
        return angle
    if angle < 0:
        return 360 + angle


def PlotSXM(sxmpath):
    try:
        raw_data = nap.read.Scan(sxmpath)
        direction = raw_data.header["scan_dir"]
        origin = "lower" if direction == "up" else "upper"
        scan_range = raw_data.header["scan_range"]
        scan_offset = raw_data.header["scan_offset"]
        scan_angle = raw_data.header["scan_angle"]
        topo = raw_data.signals["Z"]["forward"]
        topo = subtractMeanPlane(topo)
        topo = img_rotate_for_box(topo, eval(scan_angle))
        ax.imshow(
            topo,
            origin=origin,
            cmap=topo_colormap("Blues_r"),
            extent=(
                (scan_offset[0] - scan_range[0] / 2) * 1e9,
                (scan_offset[0] + scan_range[0] / 2) * 1e9,
                (scan_offset[1] - scan_range[1] / 2) * 1e9,
                (scan_offset[1] + scan_range[1] / 2) * 1e9,
            ),
        )
    except Exception:
        return


# 开始做ppt
if __name__ == "__main__":
    for k in range(len(DataFolderpath)):
        Folderpath = DataFolderpath[k]
        print(f"正在处理{Folderpath}")
        # 新建一张ppt
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 插入文本框
        text_box = slide.shapes.add_textbox(Cm(4), Cm(5.5), Inches(4), Inches(1))
        # 在文字框中添加文字
        parts = Folderpath.parts
        name_folder = "/".join(parts[-3:])

        time_folder = str(datetime.fromtimestamp(get_creation_time(Folderpath)))
        time_folder = time_folder.split(".")[0]
        text_frame = text_box.text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.text = name_folder
        paragraph.font.size = Pt(44)
        # 添加分割线
        line = slide.shapes.add_shape(
            autoshape_type_id=1,  # 形状类型为直线
            left=Cm(0),
            top=Cm(9.5),
            width=Cm(25.4),
            height=Cm(1.3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 255)
        # 添加创建时间
        text_box = slide.shapes.add_textbox(Cm(4), Cm(13), Inches(4), Inches(1))
        text_frame = text_box.text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.text = time_folder
        paragraph.font.size = Pt(18)

        # 读出所有文件名
        ReadAllFiles(Folderpath)
        Files = list(filter(None, Files))
        # 按照数据类型分类
        SortFiles(Files)
        # 将分好的数据组按照时间排序
        DatFiles = sort_files_by_creation_time(DatFiles)
        SxmFiles = sort_files_by_creation_time(SxmFiles)
        LinecutFiles = sort_files_by_creation_time(LinecutFiles)
        MapFiles = sort_files_by_creation_time(MapFiles)

        # sxm扩展中会有map中导出的sxm类型文件,这种文件无法读取需要手动删掉
        i = 0
        while i < len(SxmFiles):
            try:
                SXM(SxmFiles[i])
                i = i + 1
            except ValueError:
                del SxmFiles[i]
                continue
        # 画出形貌数据并放在ppt中
        if len(SxmFiles) > 0:
            for i in range(len(SxmFiles)):
                if i % 10 == 0:
                    # 新建一张ppt
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # 添加PPT页标题
                    text_box = slide.shapes.add_textbox(
                        Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                    )
                    text_frame = text_box.text_frame
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = name_folder
                    paragraph.font.size = Pt(24)
                    # 添加分割线
                    line = slide.shapes.add_shape(
                        autoshape_type_id=1,  # 形状类型为直线
                        left=Cm(0),
                        top=Cm(1.33),
                        width=Cm(25.4),
                        height=Cm(0.25),
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                    # 插入文本框
                    left = Inches(0.2)
                    top = Inches(0.6)
                    width = Inches(4)
                    height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    # 在文字框中添加文字
                    text_frame = text_box.text_frame
                    text_frame.text = "Topography"
                    fig0 = plt.figure(figsize=(1.2, 1.2))
                    ax = fig0.add_subplot(111)
                    ax.tick_params(labelsize=2)
                    ax.set_xlim(-400, 400)
                    ax.set_ylim(-400, 400)

                topopath = SxmFiles[i]
                try:
                    SXM(topopath)
                    PlotSXM(topopath)
                    raw_data = nap.read.Scan(topopath)
                except Exception as e:
                    print(f"Warning: failed to read SXM file {topopath}, skipping: {e}")
                    continue
                setpointI = (
                    float(raw_data.header["z-controller"]["Setpoint"][0].split(" ")[0])
                    * 1e12
                )
                scan_range = raw_data.header["scan_range"]
                angle = float(raw_data.header["scan_angle"])
                time = str(datetime.fromtimestamp(get_creation_time(topopath)))
                time = time.split(".")[0]
                parts = Path(topopath).parts
                name = "/".join(parts[-3:])
                divider = 1
                if "d1" in str(topopath):
                    divider = 1
                if "d10" in str(topopath):
                    divider = 10
                if "d100" in str(topopath):
                    divider = 100
                setpointV = raw_data.header["bias"] * 1000 / divider
                # 关于ppt中旋转角度的确定
                direction = raw_data.header["scan_dir"]
                # Rotation sign must follow the scan direction: down scans use
                # -angle, matching the marker-plot convention used elsewhere.
                rotation_angle = -angle if direction != "up" else angle
                angle_ppt = angle_def(rotation_angle)

                if 0 <= (i % 10) <= 4:
                    xi = -0.1 + 5.1 * (i % 10)
                    yi = 4.8
                    left = Cm(xi)
                    top = Cm(yi - 2.8)
                    width = Inches(1)
                    height = Inches(0.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    # 设置文本框的字体大小
                    paragraph = text_frame.add_paragraph()

                    paragraph.text = (
                        name
                        + " \n"
                        + str(round(scan_range[0] * 1e9, 1))
                        + "nm"
                        + "×"
                        + str(round(scan_range[1] * 1e9, 1))
                        + "nm"
                        + "  "
                        + str(round(angle))
                        + "°"
                        + direction
                        + "\n"
                        + "Vs="
                        + str(round(setpointV, 2))
                        + "mV"
                        + "  It="
                        + str(round(setpointI))
                        + "pA\n"
                        + time
                    )

                    paragraph.font.size = Pt(11)
                    picture = slide.shapes.add_picture(
                        str(Storagepath / "temp.tif"), left=Cm(xi), top=Cm(yi)
                    )
                    picture.rotation = angle_ppt

                if 5 <= (i % 10) <= 9:
                    xi = -0.1 + 5.1 * (i % 10 - 5)
                    yi = 12.0
                    left = Cm(xi)
                    top = Cm(yi - 2.8)
                    width = Inches(1)
                    height = Inches(0.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    # 设置文本框的字体大小
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = (
                        name
                        + " \n"
                        + str(round(scan_range[0] * 1e9, 1))
                        + "nm"
                        + "×"
                        + str(round(scan_range[1] * 1e9, 1))
                        + "nm"
                        + "  "
                        + str(round(angle))
                        + "°"
                        + direction
                        + "\n"
                        + "Vs="
                        + str(round(setpointV, 0))
                        + "mV"
                        + "  It="
                        + str(round(setpointI, 0))
                        + "pA\n"
                        + time
                    )
                    paragraph.font.size = Pt(11)
                    picture = slide.shapes.add_picture(
                        str(Storagepath / "temp.tif"), left=Cm(xi), top=Cm(yi)
                    )
                    picture.rotation = angle_ppt

                if i % 10 == 9 or i == (len(SxmFiles) - 1):
                    fig0.savefig(
                        str(Storagepath / "boxtemp.tif"),
                        format="tif",
                        bbox_inches="tight",
                        transparent=True,
                        pad_inches=0,
                    )
                    xi = 22
                    yi = 0
                    picture = slide.shapes.add_picture(
                        str(Storagepath / "boxtemp.tif"), left=Cm(xi), top=Cm(yi)
                    )
                    picture.rotation = angle_ppt
        else:
            pass

        # 画出单谱在ppt中,并且对应在其下面画出对应的形貌,标注打谱点的位置
        if len(DatFiles) > 0:
            STS(
                DatFiles[0]
            )  # 不知道为什么循环中的第一张图坐标大小改不回来,先运行一次就能解决了
            for i in range(len(DatFiles)):
                stspath = DatFiles[i]
                STS(stspath)
                # 获取单谱的文字信息
                raw_data_sts = nap.read.Spec(stspath)
                try:
                    setpointI_sts = raw_data_sts.signals["Current [AVG] (A)"][0] * 1e12
                except Exception:
                    try:
                        setpointI_sts = raw_data_sts.signals["Current (A)"][0] * 1e12
                    except Exception:
                        setpointI_sts = 404
                divider = 1
                if "d1" in str(stspath):
                    divider = 1
                if "d10" in str(stspath):
                    divider = 10
                if "d100" in str(stspath):
                    divider = 100
                setpointV_sts = raw_data_sts.signals["Bias calc (V)"][0] * 1000 / divider
                time_sts = str(datetime.fromtimestamp(get_creation_time(stspath)))
                time_sts = time_sts.split(".")[0]
                parts = Path(stspath).parts
                name_sts = "/".join(parts[-3:])
                topopath = find_nearest_file(stspath, SxmFiles)
                if topopath != "Topography Not Found":
                    # 获取形貌的文字信息
                    raw_data_topo = nap.read.Scan(topopath)
                    setpointI_topo = (
                        float(
                            raw_data_topo.header["z-controller"]["Setpoint"][0].split(" ")[
                                0
                            ]
                        )
                        * 1e12
                    )
                    scan_range = raw_data_topo.header["scan_range"]
                    angle = float(raw_data_topo.header["scan_angle"])
                    angle_ppt = angle_def(angle)
                    time_topo = str(datetime.fromtimestamp(get_creation_time(topopath)))
                    time_topo = time_topo.split(".")[0]
                    topo_parts = Path(topopath).parts
                    name_topo = "/".join(topo_parts[-3:])
                    divider = 1
                    if "d1" in str(topopath):
                        divider = 1
                    if "d10" in str(topopath):
                        divider = 10
                    if "d100" in str(topopath):
                        divider = 100
                    setpointV_topo = raw_data_topo.header["bias"] * 1000 / divider
                    # 再画一次没有标识的形貌
                    SXM(topopath)
                # 开始将获得的照片画在ppt上并标注文字信息
                if i % 5 == 0:
                    # 新建一张ppt
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # 添加PPT页标题
                    text_box = slide.shapes.add_textbox(
                        Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                    )
                    text_frame = text_box.text_frame
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = name_folder
                    paragraph.font.size = Pt(24)
                    # 添加分割线
                    line = slide.shapes.add_shape(
                        autoshape_type_id=1,  # 形状类型为直线
                        left=Cm(0),
                        top=Cm(1.33),
                        width=Cm(25.4),
                        height=Cm(0.25),
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                    # 插入文本框
                    left = Inches(0.2)
                    top = Inches(0.6)
                    width = Inches(4)
                    height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    # 在文字框中添加文字
                    text_frame = text_box.text_frame
                    text_frame.text = "STS and its point in topography"
                # 画出单谱
                xi = -0.1 + 5.1 * (i % 5)
                yi = 4.8
                left = Cm(xi + 0.8)
                top = Cm(yi - 2.5)
                width = Inches(1)
                height = Inches(0.5)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                # 设置文本框的字体大小
                paragraph = text_frame.add_paragraph()
                paragraph.text = (
                    name_sts
                    + " \n"
                    + "It="
                    + str(round(setpointI_sts, 0))
                    + "pA\n"
                    + time_sts
                )
                paragraph.font.size = Pt(11)
                slide.shapes.add_picture(
                    str(Storagepath / "temp_sts.tif"), left=Cm(xi), top=Cm(yi)
                )
                if topopath != "Topography Not Found":
                    # 画出对应形貌
                    xi = -0.1 + 5.1 * (i % 5)
                    yi = 12.8
                    left = Cm(xi)
                    top = Cm(yi - 2.8)
                    width = Inches(1)
                    height = Inches(0.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    # 设置文本框的字体大小
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = (
                        name_topo
                        + " \n"
                        + str(round(scan_range[0] * 1e9, 1))
                        + "nm"
                        + "×"
                        + str(round(scan_range[1] * 1e9, 1))
                        + "nm"
                        + "  "
                        + str(round(angle))
                        + "°\n"
                        "Vs="
                        + str(round(setpointV_topo, 0))
                        + "mV"
                        + "  It="
                        + str(round(setpointI_topo))
                        + "pA\n"
                        + time_topo
                    )
                    paragraph.font.size = Pt(11)
                    picture = slide.shapes.add_picture(
                        str(Storagepath / "temp.tif"), left=Cm(xi), top=Cm(yi)
                    )
                    picture.rotation = angle_ppt
                    slide.shapes.add_picture(
                        str(Storagepath / "temp_ststopo.tif"), left=Cm(xi), top=Cm(yi)
                    )
                if topopath == "Topography Not Found":
                    xi = -0.1 + 5.1 * (i % 5)
                    yi = 12.8
                    left = Cm(xi)
                    top = Cm(yi)
                    width = Inches(1)
                    height = Inches(0.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    # 设置文本框的字体大小
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = topopath
                    paragraph.font.size = Pt(15)
        else:
            pass

        # 画出linecut图以及相应形貌图在ppt中
        if len(LinecutFiles):
            Linecut(LinecutFiles[0])  # 先画一次,第一次跑函数字体有bug
            for i in range(len(LinecutFiles)):
                lcpath = LinecutFiles[i]
                Linecut(lcpath)
                # 获取linecut的文字信息
                raw_data_lc = nap.read.Grid(lcpath)
                try:
                    setpointI_lc = raw_data_lc.signals["Current [AVG] (A)"][0][0][0] * 1e12
                except Exception:
                    try:
                        setpointI_lc = raw_data_lc.signals["Current (A)"][0][0][0] * 1e12
                    except Exception:
                        setpointI_lc = 404
                divider = 1
                if "d1" in str(lcpath):
                    divider = 1
                if "d10" in str(lcpath):
                    divider = 10
                if "d100" in str(lcpath):
                    divider = 100
                setpointV_lc = raw_data_lc.signals["sweep_signal"][0] * 1000 / divider
                time_lc = str(datetime.fromtimestamp(get_creation_time(lcpath)))
                time_lc = time_lc.split(".")[0]
                parts = Path(lcpath).parts
                name_lc = "/".join(parts[-3:])
                linecut_size = raw_data_lc.header["size_xy"][0] * 1e9
                angle_lc = raw_data_lc.header["angle"]
                # 获取形貌的文字信息
                topopath = find_nearest_file(lcpath, SxmFiles)
                if topopath != "Topography Not Found":
                    raw_data_topo = nap.read.Scan(topopath)
                    setpointI_topo = (
                        float(
                            raw_data_topo.header["z-controller"]["Setpoint"][0].split(" ")[
                                0
                            ]
                        )
                        * 1e12
                    )
                    scan_range = raw_data_topo.header["scan_range"]
                    angle = float(raw_data_topo.header["scan_angle"])
                    angle_ppt = angle_def(angle)
                    time_topo = str(datetime.fromtimestamp(get_creation_time(topopath)))
                    time_topo = time_topo.split(".")[0]
                    topo_parts = Path(topopath).parts
                    name_topo = "/".join(topo_parts[-3:])
                    divider = 1
                    if "d1" in str(topopath):
                        divider = 1
                    if "d10" in str(topopath):
                        divider = 10
                    if "d100" in str(topopath):
                        divider = 100
                    setpointV_topo = raw_data_topo.header["bias"] * 1000 / divider
                    # 再画一次旋转没有标记的形貌
                    SXM(topopath)
                # 开始将获得的照片画在ppt上并标注文字信息
                if i % 2 == 0:
                    # 新建一张ppt
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # 添加PPT页标题
                    text_box = slide.shapes.add_textbox(
                        Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                    )
                    text_frame = text_box.text_frame
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = name_folder
                    paragraph.font.size = Pt(24)
                    # 添加分割线
                    line = slide.shapes.add_shape(
                        autoshape_type_id=1,  # 形状类型为直线
                        left=Cm(0),
                        top=Cm(1.18),
                        width=Cm(25.4),
                        height=Cm(0.25),
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                    # 插入文本框
                    left = Inches(0)
                    top = Cm(1.26)
                    width = Inches(4)
                    height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    # 在文字框中添加文字
                    text_frame = text_box.text_frame
                    text_frame.text = "Linecut and its location on topography:"
                # 画出linecut
                xi = -0.1 + 13 * (i % 2)
                yi = 3.8
                left = Cm(xi + 3)
                top = Cm(yi - 2.5)
                width = Inches(1)
                height = Inches(0.5)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                # 设置linecut文本框的字体
                paragraph = text_frame.add_paragraph()
                paragraph.text = (
                    name_lc
                    + "\n"
                    + "Length="
                    + str(round(linecut_size, 1))
                    + "nm"
                    + "  Angle="
                    + str(round(angle_lc))
                    + "°\n"
                    "It=" + str(round(setpointI_lc, 0)) + "pA\n" + time_lc
                )
                paragraph.font.size = Pt(11)
                slide.shapes.add_picture(
                    str(Storagepath / "temp_lc.tif"), left=Cm(xi), top=Cm(yi)
                )
                # 画出对应形貌
                xi = 6.7 + 13 * (i % 2)
                yi = 12.76
                if topopath != "Topography Not Found":
                    left = Cm(xi)
                    top = Cm(yi - 2.8)
                    width = Inches(1)
                    height = Inches(0.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    # 设置文本框的字体
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = (
                        name_topo
                        + " \n"
                        + str(round(scan_range[0] * 1e9, 1))
                        + "nm"
                        + "×"
                        + str(round(scan_range[1] * 1e9, 1))
                        + "nm"
                        + "  "
                        + str(round(angle))
                        + "°\n"
                        "Vs="
                        + str(round(setpointV_topo, 1))
                        + "mV"
                        + "  It="
                        + str(round(setpointI_topo))
                        + "pA\n"
                        + time_topo
                    )
                    paragraph.font.size = Pt(11)
                    picture = slide.shapes.add_picture(
                        str(Storagepath / "temp.tif"), left=Cm(xi), top=Cm(yi)
                    )
                    picture.rotation = angle_ppt
                    slide.shapes.add_picture(
                        str(Storagepath / "temp_lctopo.tif"), left=Cm(xi), top=Cm(yi)
                    )
                # 画出overlap
                slide.shapes.add_picture(
                    str(Storagepath / "temp_ol.tif"), left=Cm(xi - 7), top=Cm(yi - 0.5)
                )
        else:
            pass

        # 画出Map数据并放在ppt中
        if len(MapFiles) > 0:
            for j in range(len(MapFiles)):
                mappath = MapFiles[j]
                topopath = find_nearest_file(mappath, SxmFiles)
                raw_data = nap.read.Grid(mappath)
                topo_inmap = raw_data.signals["topo"]
                try:
                    setpoint = raw_data.signals["Current (A)"][0][0][0] * 1e12
                except Exception:
                    setpoint = raw_data.signals["Current [AVG] (A)"][0][0][0] * 1e12
                bias = (
                    raw_data.signals["sweep_signal"] * 1000
                )  # 这里只是用来计数,不用管divider的换算问题
                scan_range = raw_data.header["size_xy"]
                angle = float(raw_data.header["angle"])
                time = str(datetime.fromtimestamp(get_creation_time(mappath)))
                time = time.split(".")[0]
                parts = Path(mappath).parts
                name = "/".join(parts[-3:])
                if topopath != "Topography Not Found":
                    raw_data_topo = nap.read.Scan(topopath)
                    scan_range_topo = raw_data_topo.header["scan_range"]
                    angle_topo = float(raw_data_topo.header["scan_angle"])
                    time_topo = str(datetime.fromtimestamp(get_creation_time(topopath)))
                    time_topo = time_topo.split(".")[0]
                    topo_parts = Path(topopath).parts
                    name_topo = "/".join(topo_parts[-3:])
                    angle_topo_ppt = angle_def(angle_topo)
                # 画出从map中获取的形貌图
                # map和形貌给出的角度相同,但所有的map都是从下往上扫,注意angle和origin
                fig = plt.figure(figsize=(2.55, 2.55))
                ax0 = fig.add_subplot(111)
                plt.imshow(
                    topo_inmap,
                    origin="lower",
                    cmap=topo_colormap("Blues_r"),
                    extent=(0, scan_range[0] * 1e9, 0, scan_range[1] * 1e9),
                )
                plt.xticks([])
                plt.yticks([])
                plt.axis("off")
                plt.savefig(
                    str(Storagepath / "temp_inmap.tif"),
                    format="tif",
                    bbox_inches="tight",
                    transparent=True,
                    pad_inches=0,
                )
                plt.close()
                i = 0
                # 由于要画动图,逻辑和前边不一样,先将所有map储存在一个folder_map文件夹里
                (Storagepath / "folder_map").mkdir(parents=True, exist_ok=True)
                for n in range(len(bias)):
                    ShowMap(mappath, n)
                # 创建一个动画,将多张图片合成为一个动画
                fig, ax = plt.subplots(figsize=(2, 2))
                ax.axis("off")  # 不显示坐标轴
                num_frames = len(bias)
                # 创建动画
                animation = FuncAnimation(
                    fig, update_map, frames=range(num_frames), interval=100
                )
                # 如果需要保存为文件
                animation.save(
                    str(Storagepath / "animation_map.gif"), writer="pillow", fps=4
                )
                plt.close(fig)  # close the animation figure to prevent leaks
                angle_ppt = angle_def(angle)
                if topopath != "Topography Not Found":
                    direction = raw_data_topo.header["scan_dir"]
                for i in range(len(bias)):
                    if i % 10 == 0:
                        # 新建一张ppt
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        # 添加PPT页标题
                        text_box = slide.shapes.add_textbox(
                            Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                        )
                        text_frame = text_box.text_frame
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = name_folder
                        paragraph.font.size = Pt(24)
                        # 添加分割线
                        line = slide.shapes.add_shape(
                            autoshape_type_id=1,  # 形状类型为直线
                            left=Cm(0),
                            top=Cm(1.22),
                            width=Cm(25.4),
                            height=Cm(0.25),
                        )
                        line.fill.solid()
                        line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                        # 插入文本框
                        left = Cm(0.5)
                        top = Cm(2.2)
                        width = Inches(4)
                        height = Inches(1)
                        text_box = slide.shapes.add_textbox(left, top, width, height)
                        # 在文字框中添加文字
                        text_frame = text_box.text_frame
                        text_frame.text = (
                            "Map:"
                            + "  part #"
                            + str(int(i / 10) + 1)
                            + "/"
                            + str(int((len(bias) - 1) / 10) + 1)
                            + "\n"
                            + name
                            + "\n"
                            + str(round(scan_range[0] * 1e9, 1))
                            + "nm"
                            + "×"
                            + str(round(scan_range[1] * 1e9, 1))
                            + "nm"
                            + "  "
                            + str(round(angle))
                            + "°\n"
                            + "It="
                            + str(round(setpoint))
                            + "pA\n"
                            + time
                        )
                    # 给出其形貌图和来自map中的形貌以及gif动图
                    if i == 0:
                        picture = slide.shapes.add_picture(
                            str(Storagepath / "animation_map.gif"),
                            left=Cm(18.8),
                            top=Cm(0.8),
                        )
                        picture.rotation = angle_ppt
                        if topopath != "Topography Not Found":
                            SXM(topopath)
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "temp.tif"), left=Cm(9.5), top=Cm(1.8)
                            )
                            picture.rotation = angle_topo_ppt
                            # 插入文本框
                            left2 = Cm(9.5)
                            top2 = Cm(6.3)
                            width2 = Inches(4)
                            height2 = Inches(1)
                            text_box = slide.shapes.add_textbox(
                                left2, top2, width2, height2
                            )
                            text_frame = text_box.text_frame
                            paragraph = text_frame.add_paragraph()
                            paragraph.text = (
                                name_topo
                                + "\n"
                                + time_topo
                                + "\n"
                                + str(round(scan_range_topo[0] * 1e9, 1))
                                + "nm"
                                + "×"
                                + str(round(scan_range_topo[1] * 1e9, 1))
                                + "nm"
                                + "  "
                                + str(round(angle_topo))
                                + "°\n"
                            )
                            paragraph.font.size = Pt(11)
                        picture = slide.shapes.add_picture(
                            str(Storagepath / "temp_inmap.tif"), left=Cm(14.5), top=Cm(1.8)
                        )
                        picture.rotation = angle_ppt
                        # 插入文本框
                        left3 = Cm(15.9)
                        top3 = Cm(6.7)
                        width3 = Inches(4)
                        height3 = Inches(1)
                        text_box = slide.shapes.add_textbox(left3, top3, width3, height3)
                        text_frame = text_box.text_frame
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = "topo in map"
                        paragraph.font.size = Pt(11)
                    if 0 <= (i % 10) <= 4:
                        xi = -0.1 + 5.1 * (i % 10)
                        yi = 8.8
                        picture = slide.shapes.add_picture(
                            str(Storagepath / "folder_map" / f"temp_map_{i}.tif"),
                            left=Cm(xi),
                            top=Cm(yi),
                        )
                        picture.rotation = angle_ppt
                    if 5 <= (i % 10) <= 9:
                        xi = -0.1 + 5.1 * (i % 10 - 5)
                        yi = 14
                        picture = slide.shapes.add_picture(
                            str(Storagepath / "folder_map" / f"temp_map_{i}.tif"),
                            left=Cm(xi),
                            top=Cm(yi),
                        )
                        picture.rotation = angle_ppt
                if QPI_switch == "on":
                    (str(Storagepath / "folder_QPI")).mkdir(parents=True, exist_ok=True)
                    for n in range(len(bias)):
                        QPI(mappath, n)
                    # 创建一个动画,将多张图片合成为一个动画
                    fig, ax = plt.subplots(figsize=(2, 2))
                    ax.axis("off")  # 不显示坐标轴
                    num_frames = len(bias)
                    # 创建动画
                    animation = FuncAnimation(
                        fig, update_QPI, frames=range(num_frames), interval=100
                    )
                    # 如果需要保存为文件
                    animation.save(
                        str(Storagepath / "animation_QPI.gif"), writer="pillow", fps=4
                    )
                    plt.close(fig)  # close the animation figure to prevent leaks
                    # 贴ppt
                    for i in range(len(bias)):
                        if i % 10 == 0:
                            # 新建一张ppt
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            # 添加PPT页标题
                            text_box = slide.shapes.add_textbox(
                                Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                            )
                            text_frame = text_box.text_frame
                            paragraph = text_frame.add_paragraph()
                            paragraph.text = name_folder
                            paragraph.font.size = Pt(24)
                            # 添加分割线
                            line = slide.shapes.add_shape(
                                autoshape_type_id=1,  # 形状类型为直线
                                left=Cm(0),
                                top=Cm(1.22),
                                width=Cm(25.4),
                                height=Cm(0.25),
                            )
                            line.fill.solid()
                            line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                            # 插入文本框
                            left = Cm(0.5)
                            top = Cm(2.2)
                            width = Inches(4)
                            height = Inches(1)
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            # 在文字框中添加文字
                            text_frame = text_box.text_frame
                            text_frame.text = (
                                "QPI:"
                                + "  part #"
                                + str(int(i / 10) + 1)
                                + "/"
                                + str(int((len(bias) - 1) / 10) + 1)
                                + "\n"
                                + name
                            )
                        if i == 0:
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "animation_QPI.gif"),
                                left=Cm(13.6),
                                top=Cm(0.8),
                            )
                            picture.rotation = angle_ppt
                        # 给出其形貌图和来自map中的形貌
                        if 0 <= (i % 10) <= 4:
                            xi = -0.1 + 5.1 * (i % 10)
                            yi = 8.8
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "folder_QPI" / f"temp_QPI_{i}.tif"),
                                left=Cm(xi),
                                top=Cm(yi),
                            )
                            picture.rotation = angle_ppt
                        if 5 <= (i % 10) <= 9:
                            xi = -0.1 + 5.1 * (i % 10 - 5)
                            yi = 14
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "folder_QPI" / f"temp_QPI_{i}.tif"),
                                left=Cm(xi),
                                top=Cm(yi),
                            )
                            picture.rotation = angle_ppt
                if mapI_switch == "on":
                    (str(Storagepath / "folder_mapI")).mkdir(parents=True, exist_ok=True)
                    for n in range(len(bias)):
                        ShowMapI(mappath, n)
                    # 创建一个动画,将多张图片合成为一个动画
                    fig, ax = plt.subplots(figsize=(2, 2))
                    ax.axis("off")  # 不显示坐标轴
                    num_frames = len(bias)
                    # 创建动画
                    animation = FuncAnimation(
                        fig, update_mapI, frames=range(num_frames), interval=100
                    )
                    # 如果需要保存为文件
                    animation.save(
                        str(Storagepath / "animation_mapI.gif"), writer="pillow", fps=4
                    )
                    plt.close(fig)  # close the animation figure to prevent leaks
                    for i in range(len(bias)):
                        if i % 10 == 0:
                            # 新建一张ppt
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            # 添加PPT页标题
                            text_box = slide.shapes.add_textbox(
                                Cm(-0.1), Cm(-0.7), Inches(4), Inches(1)
                            )
                            text_frame = text_box.text_frame
                            paragraph = text_frame.add_paragraph()
                            paragraph.text = name_folder
                            paragraph.font.size = Pt(24)
                            # 添加分割线
                            line = slide.shapes.add_shape(
                                autoshape_type_id=1,  # 形状类型为直线
                                left=Cm(0),
                                top=Cm(1.22),
                                width=Cm(25.4),
                                height=Cm(0.25),
                            )
                            line.fill.solid()
                            line.fill.fore_color.rgb = RGBColor(0, 0, 255)
                            # 插入文本框
                            left = Cm(0.5)
                            top = Cm(2.2)
                            width = Inches(4)
                            height = Inches(1)
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            # 在文字框中添加文字
                            text_frame = text_box.text_frame
                            text_frame.text = (
                                "Current in Map"
                                + "  part #"
                                + str(int(i / 10) + 1)
                                + "/"
                                + str(int((len(bias) - 1) / 10) + 1)
                                + "\n"
                                + name
                            )
                        if i == 0:
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "animation_mapI.gif"),
                                left=Cm(13.6),
                                top=Cm(0.8),
                            )
                            picture.rotation = angle_ppt
                        # 给出其形貌图和来自map中的形貌
                        if 0 <= (i % 10) <= 4:
                            xi = -0.1 + 5.1 * (i % 10)
                            yi = 8.8
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "folder_mapI" / f"temp_mapI_{i}.tif"),
                                left=Cm(xi),
                                top=Cm(yi),
                            )
                            picture.rotation = angle_ppt
                        if 5 <= (i % 10) <= 9:
                            xi = -0.1 + 5.1 * (i % 10 - 5)
                            yi = 14
                            picture = slide.shapes.add_picture(
                                str(Storagepath / "folder_mapI" / f"temp_mapI_{i}.tif"),
                                left=Cm(xi),
                                top=Cm(yi),
                            )
                            picture.rotation = angle_ppt
        else:
            pass
        # 每次对一个文件夹数据处理完,都清空这些数组,以便下一个文件夹的处理
        DatFiles = []
        SxmFiles = []
        LinecutFiles = []
        MapFiles = []
        Files = []

    # 删除文件中保存的图片
    # "*temp*.tif" also matches boxtemp.tif, which "temp*.tif" missed
    for temp_file in Storagepath.glob("*temp*.tif"):
        temp_file.unlink(missing_ok=True)
    shutil.rmtree(str(Storagepath / "folder_map"), ignore_errors=True)
    (Storagepath / "temp_sts.tif").unlink(missing_ok=True)
    (Storagepath / "temp_ststopo.tif").unlink(missing_ok=True)
    (Storagepath / "temp_lc.tif").unlink(missing_ok=True)
    (Storagepath / "temp_ol.tif").unlink(missing_ok=True)
    (Storagepath / "temp_lctopo.tif").unlink(missing_ok=True)
    (Storagepath / "temp_map.tif").unlink(missing_ok=True)
    (Storagepath / "temp_inmap.tif").unlink(missing_ok=True)
    (Storagepath / "temp_QPI.tif").unlink(missing_ok=True)
    (Storagepath / "temp_mapI.tif").unlink(missing_ok=True)

    prs.save(str(Storagepath / PPTname))

    (Storagepath / "animation_map.gif").unlink(missing_ok=True)
    (Storagepath / "animation_QPI.gif").unlink(missing_ok=True)
    (Storagepath / "animation_mapI.gif").unlink(missing_ok=True)

    shutil.rmtree(str(Storagepath / "folder_map"), ignore_errors=True)
    shutil.rmtree(str(Storagepath / "folder_QPI"), ignore_errors=True)
    shutil.rmtree(str(Storagepath / "folder_mapI"), ignore_errors=True)
