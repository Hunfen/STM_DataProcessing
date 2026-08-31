#!/Users/hunfen/Documents/GitHub/STM_DataProcessing/.venv/bin/python


import logging
import platform
import shutil
import sys
from datetime import datetime
from pathlib import Path

import matplotlib
import matplotlib.pyplot as plt
import nanonispy as nap
from matplotlib.animation import FuncAnimation

# isort: off
try:
    from .plot_funcs import (
        angle_def,
        get_divider,
        img_rotate_for_box,
        plot_linecut,
        plot_map_bias,
        plot_map_current_bias,
        plot_qpi_bias,
        plot_sts,
        plot_sxm_topo,
        subtractMeanPlane,
        update_frame_from_dir,
    )
except ImportError:  # allow running as a standalone script
    from plot_funcs import (
        angle_def,
        get_divider,
        img_rotate_for_box,
        plot_linecut,
        plot_map_bias,
        plot_map_current_bias,
        plot_qpi_bias,
        plot_sts,
        plot_sxm_topo,
        subtractMeanPlane,
        update_frame_from_dir,
    )
# isort: on

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches, Pt

logging.basicConfig(
    level=logging.INFO, format="%(message)s", handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)


# ---------- File processing helper functions ----------
def get_creation_time(file_path):
    stat = Path(file_path).stat()
    return stat.st_birthtime if platform.system() == "Darwin" else stat.st_ctime


def sort_files_by_creation_time(files):
    file_times = [(f, get_creation_time(f)) for f in files]
    sorted_files = sorted(file_times, key=lambda x: x[1])
    return [f[0] for f in sorted_files]


def read_all_files(folder_path: Path):
    return [file for file in folder_path.rglob("*") if file.is_file()]


def sort_files(files):
    dat_files, sxm_files, linecut_files, map_files = [], [], [], []
    for each in files:
        if each.suffix == ".dat":
            try:
                nap.read.Spec(str(each)).signals["Bias calc (V)"]
                dat_files.append(each)
            except Exception:
                pass
        elif each.suffix == ".sxm":
            sxm_files.append(each)
        elif each.suffix == ".3ds":
            raw = nap.read.Grid(str(each))
            try:
                data = raw.signals["LI Demod 1 X (A)"]
            except Exception:
                data = raw.signals["LI Demod 1 X [AVG] (A)"]
            if data.shape[0] == 1:
                linecut_files.append(each)
            else:
                map_files.append(each)
    return dat_files, sxm_files, linecut_files, map_files


def find_nearest_file(target_path, file_list):
    if not file_list:
        return "Topography Not Found"
    target_time = get_creation_time(target_path)
    nearest = None
    min_diff = float("inf")
    for f in file_list:
        diff = target_time - get_creation_time(f)
        if 0 < diff < min_diff:
            min_diff = diff
            nearest = f
    return nearest if nearest else "Topography Not Found"


# ---------- PPT layout helper functions ----------
def add_title_slide(slide, folder_name, folder_time):
    text_box = slide.shapes.add_textbox(Cm(4), Cm(5.5), Inches(4), Inches(1))
    paragraph = text_box.text_frame.add_paragraph()
    paragraph.text = folder_name
    paragraph.font.size = Pt(44)

    line = slide.shapes.add_shape(1, Cm(0), Cm(9.5), Cm(25.4), Cm(1.3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 255)

    text_box = slide.shapes.add_textbox(Cm(4), Cm(13), Inches(4), Inches(1))
    paragraph = text_box.text_frame.add_paragraph()
    paragraph.text = folder_time.split(".")[0]
    paragraph.font.size = Pt(18)


def add_section_header(slide, title):
    text_box = slide.shapes.add_textbox(Cm(-0.1), Cm(-0.7), Inches(4), Inches(1))
    paragraph = text_box.text_frame.add_paragraph()
    paragraph.text = title
    paragraph.font.size = Pt(24)

    line = slide.shapes.add_shape(1, Cm(0), Cm(1.33), Cm(25.4), Cm(0.25))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 255)


def main() -> None:
    """Interactive entry point: gather settings, build and save the PPT.

    Prompts for data folder paths, PPT storage path/name and image
    switches, then processes every folder and saves the .pptx file."""

    # ================= Interactive Input Section =================
    logger.info(
        "Please enter data processing folder path (separate multiple paths with commas):"
    )
    input_paths = input("> ").strip()
    if input_paths:
        DataFolderpath = [Path(p.strip()) for p in input_paths.split(",")]
    else:
        sys.exit("Error: No data folder path entered, script terminated.")

    logger.info(
        "\nPlease enter PPT storage directory path (press Enter to save in data folder path):"
    )
    input_storage = input("> ").strip()
    if input_storage:
        Storagepath = Path(input_storage)
    else:
        Storagepath = Path(DataFolderpath[0])
        logger.info(
            f"Output path not specified, will save PPT in data folder: {Storagepath}"
        )

    logger.info(
        "\nPlease enter PPT file name (without path, press Enter to auto-generate):"
    )
    input_pptname = input("> ").strip()
    if input_pptname:
        if not input_pptname.endswith(".pptx"):
            input_pptname += ".pptx"
        PPTname = input_pptname
    else:
        PPTname = f"auto_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        logger.info(f"Auto-generated file name: {PPTname}")

    logger.info(
        "\nOutput QPI images? Enter 'on' or 'off' (press Enter for default 'on'):"
    )
    QPI_switch = input("> ").strip().lower() in ("", "on")
    logger.info(f"QPI output: {'on' if QPI_switch else 'off'}")

    logger.info(
        "\nOutput current maps? Enter 'on' or 'off' (press Enter for default 'off'):"
    )
    mapI_switch = input("> ").strip().lower() == "on"
    logger.info(f"Map current output: {'on' if mapI_switch else 'off'}")

    logger.info(
        "\nApply Gaussian smoothing to spectra/maps? Enter 'on' or 'off' (press Enter for default 'off'):"
    )
    smooth_switch = input("> ").strip().lower() == "on"
    logger.info(f"Smoothing: {'on' if smooth_switch else 'off'}")

    prs = Presentation()
    matplotlib.use("Agg")
    plt.rcParams["font.sans-serif"] = ["Arial"]
    plt.rcParams["axes.unicode_minus"] = True
    plt.rcParams.update({"font.size": 22})

    # ---------- Main process ----------
    for _, Folderpath in enumerate(DataFolderpath):
        logger.info(f"Processing {Folderpath}")
        # First title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        name_folder = "/".join(Folderpath.parts[-3:])
        time_folder = str(datetime.fromtimestamp(get_creation_time(Folderpath)))
        add_title_slide(slide, name_folder, time_folder)

        # File classification and sorting
        all_files = read_all_files(Folderpath)
        DatFiles, SxmFiles, LinecutFiles, MapFiles = sort_files(all_files)
        DatFiles = sort_files_by_creation_time(DatFiles)
        SxmFiles = sort_files_by_creation_time(SxmFiles)
        LinecutFiles = sort_files_by_creation_time(LinecutFiles)
        MapFiles = sort_files_by_creation_time(MapFiles)

        # Clean up invalid SXM files
        valid_sxm = []
        for f in SxmFiles:
            try:
                plot_sxm_topo(f, Storagepath / "temp_test.tif")
                valid_sxm.append(f)
            except Exception:
                pass
        SxmFiles = valid_sxm

        # ---- topos ----
        if SxmFiles:
            fig0 = None
            for i, topo_path in enumerate(SxmFiles):
                if i % 8 == 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_section_header(slide, name_folder)
                    fig0 = plt.figure(figsize=(1.2, 1.2))
                    ax0 = fig0.add_subplot(111)
                    ax0.tick_params(labelsize=2)
                    ax0.set_xlim(-400, 400)
                    ax0.set_ylim(-400, 400)

                plot_sxm_topo(topo_path, Storagepath / "temp.tif")
                raw_data = nap.read.Scan(str(topo_path))
                setpointI = (
                    float(raw_data.header["z-controller"]["Setpoint"][0].split(" ")[0])
                    * 1e12
                )
                scan_range = raw_data.header["scan_range"]
                angle = float(raw_data.header["scan_angle"])
                time_str = str(
                    datetime.fromtimestamp(get_creation_time(topo_path))
                ).split(".")[0]
                # name = "/".join(Path(topo_path).parts[-3:])
                name = Path(topo_path).name
                divider = get_divider(topo_path)
                setpointV = raw_data.header["bias"] * 1000 / divider
                direction = raw_data.header["scan_dir"]
                angle_ppt = angle_def(angle)

                # Use topo data for overview box
                topo = raw_data.signals["Z"]["forward"]
                topo_plane = subtractMeanPlane(topo)
                topo_rotated = img_rotate_for_box(
                    topo_plane, float(raw_data.header["scan_angle"])
                )
                ax0.imshow(
                    topo_rotated,
                    origin="lower" if direction == "up" else "upper",
                    cmap="Blues_r",
                    extent=(
                        (float(raw_data.header["scan_offset"][0]) - scan_range[0] / 2)
                        * 1e9,
                        (float(raw_data.header["scan_offset"][0]) + scan_range[0] / 2)
                        * 1e9,
                        (float(raw_data.header["scan_offset"][1]) - scan_range[1] / 2)
                        * 1e9,
                        (float(raw_data.header["scan_offset"][1]) + scan_range[1] / 2)
                        * 1e9,
                    ),
                )

                row = i % 8
                col = row % 4
                if row < 4:
                    xi = 6.5 * col  # 上方行
                    yi = 2.0
                else:
                    xi = 6.5 * col  # 下方行
                    yi = 11.0

                pic = slide.shapes.add_picture(
                    str(Storagepath / "temp.tif"), Cm(xi), Cm(yi)
                )
                pic.rotation = angle_ppt

                txt_top = pic.top + pic.height + Cm(0.2)
                txt = slide.shapes.add_textbox(Cm(xi), txt_top, Cm(6.5), Cm(2.64))
                txt.text_frame.clear()
                txt.text_frame.paragraphs[0].text = (
                    f"{time_str}\n"
                    f"{name}\n"
                    f"{round(scan_range[0] * 1e9, 1)}×{round(scan_range[1] * 1e9, 1)} nm²  {round(angle)}°{direction}\n"
                    f"{setpointV:.2f} mV, {setpointI:.0f} pA"
                )
                txt.text_frame.paragraphs[0].font.size = Pt(14)

                if (i + 1) % 8 == 0 or i == len(
                    SxmFiles
                ) - 1:  # 每页结束或最后一张图时保存概述图
                    fig0.savefig(
                        str(Storagepath / "boxtemp.tif"),
                        format="tif",
                        bbox_inches="tight",
                        transparent=True,
                        pad_inches=0,
                    )
                    pic = slide.shapes.add_picture(
                        str(Storagepath / "boxtemp.tif"), Cm(22), Cm(0)
                    )
                    pic.rotation = angle_ppt

        # ---- spectrum ----
        if DatFiles:
            for i, sts_path in enumerate(DatFiles):
                if i == 0:  # Warm-up
                    plot_sts(sts_path, None, Storagepath, smooth_switch)
                topopath = find_nearest_file(sts_path, SxmFiles)
                topo_for_sts = topopath if topopath != "Topography Not Found" else None
                sts_img, topo_marked_img = plot_sts(
                    sts_path, topo_for_sts, Storagepath, smooth_switch
                )

                if i % 5 == 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_section_header(slide, name_folder)

                xi = -0.1 + 5.1 * (i % 5)
                yi = 4.8
                slide.shapes.add_picture(str(sts_img), Cm(xi), Cm(yi))
                # txt = slide.shapes.add_textbox(
                #     Cm(xi + 0.8), Cm(yi - 2.5), Inches(1), Inches(0.5)
                # )
                # txt.text_frame.add_paragraph().text = (
                #     f"{name_sts}\nIt={round(setpointI_sts)}pA\n{time_sts}"
                # )
                # txt.text_frame.paragraphs[0].font.size = Pt(11)

                if topo_marked_img:
                    xi = -0.1 + 5.1 * (i % 5)
                    yi = 12.8
                    slide.shapes.add_picture(str(topo_marked_img), Cm(xi), Cm(yi))

                else:
                    if topopath == "Topography Not Found":
                        slide.shapes.add_textbox(
                            Cm(xi), Cm(12.8), Inches(1), Inches(0.5)
                        ).text = "Topography Not Found"

        # ----  linecut ----
        if LinecutFiles:
            for i, lc_path in enumerate(LinecutFiles):
                if i == 0:  # Warm-up
                    plot_linecut(lc_path, None, Storagepath, smooth_switch)
                topopath = find_nearest_file(lc_path, SxmFiles)
                topo_for_lc = topopath if topopath != "Topography Not Found" else None
                lc_img, ol_img, topo_marked_img = plot_linecut(
                    lc_path, topo_for_lc, Storagepath, smooth_switch
                )

                raw_lc = nap.read.Grid(str(lc_path))
                try:
                    current_A = raw_lc.signals["Current [AVG] (A)"][0][0][0]
                except Exception:
                    current_A = raw_lc.signals["Current (A)"][0][0][0]
                setpointI_lc = current_A * 1e12
                time_lc = str(datetime.fromtimestamp(get_creation_time(lc_path))).split(
                    "."
                )[0]
                name_lc = "/".join(Path(lc_path).parts[-3:])
                lc_size = raw_lc.header["size_xy"][0] * 1e9
                angle_lc = raw_lc.header["angle"]

                if i % 2 == 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_section_header(slide, name_folder)

                xi = -0.1 + 13 * (i % 2)
                yi = 3.8
                slide.shapes.add_picture(str(lc_img), Cm(xi), Cm(yi))
                txt = slide.shapes.add_textbox(
                    Cm(xi + 3), Cm(yi - 2.5), Inches(1), Inches(0.5)
                )
                txt.text_frame.add_paragraph().text = f"{name_lc}\nLength={round(lc_size, 1)}nm  Angle={round(angle_lc)}°\nIt={round(setpointI_lc)}pA\n{time_lc}"
                txt.text_frame.paragraphs[0].font.size = Pt(11)

                xi = 6.7 + 13 * (i % 2)
                yi = 12.76
                if topo_marked_img:
                    slide.shapes.add_picture(str(topo_marked_img), Cm(xi), Cm(yi))
                slide.shapes.add_picture(str(ol_img), Cm(xi - 7), Cm(yi - 0.5))

        # ----  Map data processing (including QPI / current maps) ----
        if MapFiles:
            for _, mappath in enumerate(MapFiles):
                raw_map = nap.read.Grid(str(mappath))
                topopath = find_nearest_file(mappath, SxmFiles)
                if topopath != "Topography Not Found":
                    raw_topo = nap.read.Scan(str(topopath))
                    scan_range_topo = raw_topo.header["scan_range"]
                    angle_topo = float(raw_topo.header["scan_angle"])
                    angle_topo_ppt = angle_def(angle_topo)
                    time_topo = str(
                        datetime.fromtimestamp(get_creation_time(topopath))
                    ).split(".")[0]
                    name_topo = "/".join(Path(topopath).parts[-3:])
                else:
                    raw_topo = None

                # Get basic parameters
                try:
                    setpoint = raw_map.signals["Current (A)"][0][0][0] * 1e12
                except Exception:
                    setpoint = raw_map.signals["Current [AVG] (A)"][0][0][0] * 1e12
                bias_signal = raw_map.signals["sweep_signal"] * 1000
                n_bias = len(bias_signal)
                scan_range = raw_map.header["size_xy"]
                angle = float(raw_map.header["angle"])
                time_map = str(
                    datetime.fromtimestamp(get_creation_time(mappath))
                ).split(".")[0]
                name_map = "/".join(Path(mappath).parts[-3:])
                angle_ppt = angle_def(angle)

                # Save in-map topography image
                topo_inmap = raw_map.signals["topo"]
                fig = plt.figure(figsize=(2.55, 2.55))
                ax0 = fig.add_subplot(111)
                ax0.imshow(
                    topo_inmap,
                    origin="lower",
                    cmap="Blues_r",
                    extent=(0, scan_range[0] * 1e9, 0, scan_range[1] * 1e9),
                )
                plt.axis("off")
                plt.savefig(
                    str(Storagepath / "temp_inmap.tif"),
                    format="tif",
                    bbox_inches="tight",
                    transparent=True,
                    pad_inches=0,
                )
                plt.close()

                # Generate per-bias map images
                map_dir = Storagepath / "folder_map"
                map_dir.mkdir(parents=True, exist_ok=True)
                for n in range(n_bias):
                    plot_map_bias(mappath, n, map_dir)

                # Animated map
                fig, ax = plt.subplots(figsize=(2, 2))
                ax.axis("off")

                def anim_update(frame, map_dir=map_dir, ax=ax):
                    update_frame_from_dir(map_dir, frame, ax)

                ani = FuncAnimation(fig, anim_update, frames=n_bias, interval=100)
                ani.save(str(Storagepath / "animation_map.gif"), writer="pillow", fps=4)
                plt.close(fig)

                # Populate PPT slides
                for page_start in range(0, n_bias, 10):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_section_header(slide, name_folder)
                    txt = slide.shapes.add_textbox(
                        Cm(0.5), Cm(2.2), Inches(4), Inches(1)
                    )
                    part = page_start // 10 + 1
                    total_parts = (n_bias - 1) // 10 + 1
                    txt.text_frame.add_paragraph().text = (
                        f"Map:  part #{part}/{total_parts}\n{name_map}\n"
                        f"{round(scan_range[0] * 1e9, 1)}nm×{round(scan_range[1] * 1e9, 1)}nm  "
                        f"{round(angle)}°\nIt={round(setpoint)}pA\n{time_map}"
                    )

                    if page_start == 0:
                        pic = slide.shapes.add_picture(
                            str(Storagepath / "animation_map.gif"), Cm(18.8), Cm(0.8)
                        )
                        pic.rotation = angle_ppt
                        if raw_topo:
                            plot_sxm_topo(topopath, Storagepath / "temp.tif")
                            pic = slide.shapes.add_picture(
                                str(Storagepath / "temp.tif"), Cm(9.5), Cm(1.8)
                            )
                            pic.rotation = angle_topo_ppt
                            slide.shapes.add_textbox(
                                Cm(9.5), Cm(6.3), Inches(4), Inches(1)
                            ).text = (
                                f"{name_topo}\n{time_topo}\n"
                                f"{round(scan_range_topo[0] * 1e9, 1)}nm×{round(scan_range_topo[1] * 1e9, 1)}nm  "
                                f"{round(angle_topo)}°"
                            )
                        pic = slide.shapes.add_picture(
                            str(Storagepath / "temp_inmap.tif"), Cm(14.5), Cm(1.8)
                        )
                        pic.rotation = angle_ppt
                        slide.shapes.add_textbox(
                            Cm(15.9), Cm(6.7), Inches(4), Inches(1)
                        ).text = "topo in map"

                    # Insert map images for current page
                    for idx_offset in range(10):
                        bias_idx = page_start + idx_offset
                        if bias_idx >= n_bias:
                            break
                        if idx_offset < 5:
                            xi = -0.1 + 5.1 * idx_offset
                            yi = 8.8
                        else:
                            xi = -0.1 + 5.1 * (idx_offset - 5)
                            yi = 14
                        img_path = map_dir / f"temp_map_{bias_idx}.tif"
                        if img_path.exists():
                            pic = slide.shapes.add_picture(
                                str(img_path), Cm(xi), Cm(yi)
                            )
                            pic.rotation = angle_ppt

                # QPI processing
                if QPI_switch:
                    qpi_dir = Storagepath / "folder_QPI"
                    qpi_dir.mkdir(parents=True, exist_ok=True)
                    for n in range(n_bias):
                        plot_qpi_bias(mappath, n, qpi_dir)

                    fig, ax = plt.subplots(figsize=(2, 2))
                    ax.axis("off")

                    def anim_update_qpi(frame, qpi_dir=qpi_dir, ax=ax):
                        img_path = qpi_dir / f"temp_QPI_{frame}.tif"
                        ax.cla()
                        ax.set_axis_off()
                        img = plt.imread(str(img_path))
                        ax.imshow(img, extent=[0, img.shape[1], 0, img.shape[0]])

                    ani = FuncAnimation(
                        fig, anim_update_qpi, frames=n_bias, interval=100
                    )
                    ani.save(
                        str(Storagepath / "animation_QPI.gif"), writer="pillow", fps=4
                    )
                    plt.close(fig)

                    for page_start in range(0, n_bias, 10):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        add_section_header(slide, name_folder)
                        part = page_start // 10 + 1
                        total_parts = (n_bias - 1) // 10 + 1
                        slide.shapes.add_textbox(
                            Cm(0.5), Cm(2.2), Inches(4), Inches(1)
                        ).text = f"QPI:  part #{part}/{total_parts}\n{name_map}"
                        if page_start == 0:
                            pic = slide.shapes.add_picture(
                                str(Storagepath / "animation_QPI.gif"),
                                Cm(13.6),
                                Cm(0.8),
                            )
                            pic.rotation = angle_ppt
                        for idx_offset in range(10):
                            bias_idx = page_start + idx_offset
                            if bias_idx >= n_bias:
                                break
                            if idx_offset < 5:
                                xi = -0.1 + 5.1 * idx_offset
                                yi = 8.8
                            else:
                                xi = -0.1 + 5.1 * (idx_offset - 5)
                                yi = 14
                            img_path = qpi_dir / f"temp_QPI_{bias_idx}.tif"
                            if img_path.exists():
                                pic = slide.shapes.add_picture(
                                    str(img_path), Cm(xi), Cm(yi)
                                )
                                pic.rotation = angle_ppt

                # Current map processing
                if mapI_switch:
                    cur_dir = Storagepath / "folder_mapI"
                    cur_dir.mkdir(parents=True, exist_ok=True)
                    for n in range(n_bias):
                        plot_map_current_bias(mappath, n, cur_dir)

                    fig, ax = plt.subplots(figsize=(2, 2))
                    ax.axis("off")

                    def anim_update_cur(frame, cur_dir=cur_dir, ax=ax):
                        img_path = cur_dir / f"temp_mapI_{frame}.tif"
                        ax.cla()
                        ax.set_axis_off()
                        img = plt.imread(str(img_path))
                        ax.imshow(img, extent=[0, img.shape[1], 0, img.shape[0]])

                    ani = FuncAnimation(
                        fig, anim_update_cur, frames=n_bias, interval=100
                    )
                    ani.save(
                        str(Storagepath / "animation_mapI.gif"), writer="pillow", fps=4
                    )
                    plt.close(fig)

                    for page_start in range(0, n_bias, 10):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        add_section_header(slide, name_folder)
                        part = page_start // 10 + 1
                        total_parts = (n_bias - 1) // 10 + 1
                        slide.shapes.add_textbox(
                            Cm(0.5), Cm(2.2), Inches(4), Inches(1)
                        ).text = (
                            f"Current in Map  part #{part}/{total_parts}\n{name_map}"
                        )
                        if page_start == 0:
                            pic = slide.shapes.add_picture(
                                str(Storagepath / "animation_mapI.gif"),
                                Cm(13.6),
                                Cm(0.8),
                            )
                            pic.rotation = angle_ppt
                        for idx_offset in range(10):
                            bias_idx = page_start + idx_offset
                            if bias_idx >= n_bias:
                                break
                            if idx_offset < 5:
                                xi = -0.1 + 5.1 * idx_offset
                                yi = 8.8
                            else:
                                xi = -0.1 + 5.1 * (idx_offset - 5)
                                yi = 14
                            img_path = cur_dir / f"temp_mapI_{bias_idx}.tif"
                            if img_path.exists():
                                pic = slide.shapes.add_picture(
                                    str(img_path), Cm(xi), Cm(yi)
                                )
                                pic.rotation = angle_ppt

        # Clean up temporary files related to the folder
        for f in Storagepath.glob("*temp*.tif"):  # 改为 *temp*.tif
            f.unlink(missing_ok=True)
        shutil.rmtree(str(Storagepath / "folder_map"), ignore_errors=True)
        shutil.rmtree(str(Storagepath / "folder_QPI"), ignore_errors=True)
        shutil.rmtree(str(Storagepath / "folder_mapI"), ignore_errors=True)
        # Reset lists for next folder (can be omitted if only one folder)
        DatFiles.clear()
        SxmFiles.clear()
        LinecutFiles.clear()
        MapFiles.clear()

    # Save the PowerPoint file
    prs.save(str(Storagepath / PPTname))
    logger.info(f"PPT saved to {Storagepath / PPTname}")


if __name__ == "__main__":
    main()
